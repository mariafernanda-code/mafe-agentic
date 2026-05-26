"""Generador de presentaciones PowerPoint."""

from __future__ import annotations

import io
from pathlib import Path

import matplotlib

matplotlib.use("Agg")
import matplotlib.pyplot as plt
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt


PRIMARY = RGBColor(0x1F, 0x3A, 0x68)
ACCENT = RGBColor(0xE0, 0x6C, 0x3D)
TEXT = RGBColor(0x2A, 0x2A, 0x2A)
MUTED = RGBColor(0x70, 0x70, 0x70)
BG = RGBColor(0xF7, 0xF6, 0xF1)


def _chart_png(chart: dict) -> bytes:
    fig, ax = plt.subplots(figsize=(7.5, 4.2), dpi=160)
    labels = chart.get("labels", [])
    values = chart.get("values", [])
    title = chart.get("title", "")
    ctype = chart.get("type", "bar")

    if ctype == "bar":
        ax.bar(labels, values, color="#1F3A68")
    elif ctype == "line":
        ax.plot(labels, values, marker="o", color="#E06C3D", linewidth=2)
    elif ctype == "pie":
        ax.pie(values, labels=labels, autopct="%1.0f%%",
               colors=["#1F3A68", "#E06C3D", "#5A8AB0", "#D4A24C", "#7CA982"])
    else:
        ax.bar(labels, values, color="#1F3A68")

    if title and ctype != "pie":
        ax.set_title(title, color="#2A2A2A", fontsize=12)
    ax.spines["top"].set_visible(False)
    ax.spines["right"].set_visible(False)
    fig.tight_layout()

    buf = io.BytesIO()
    fig.savefig(buf, format="png", bbox_inches="tight")
    plt.close(fig)
    buf.seek(0)
    return buf.read()


def _txt(text_frame, text: str, *, size=18, bold=False, color=TEXT):
    text_frame.clear()
    p = text_frame.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color


def build(spec: dict, output_path: str) -> str:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    layout = prs.slide_layouts[6]

    # Portada
    slide = prs.slides.add_slide(layout)
    bg = slide.shapes.add_shape(1, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid(); bg.fill.fore_color.rgb = BG; bg.line.fill.background()

    band = slide.shapes.add_shape(1, 0, 0, Inches(0.4), prs.slide_height)
    band.fill.solid(); band.fill.fore_color.rgb = ACCENT; band.line.fill.background()

    title_box = slide.shapes.add_textbox(Inches(1.0), Inches(2.4), Inches(11.0), Inches(1.4))
    _txt(title_box.text_frame, spec.get("title", "Presentacion"),
         size=44, bold=True, color=PRIMARY)

    if spec.get("subtitle"):
        sub_box = slide.shapes.add_textbox(Inches(1.0), Inches(3.9), Inches(11.0), Inches(0.8))
        _txt(sub_box.text_frame, spec["subtitle"], size=20, color=MUTED)

    # Slides de contenido
    for sd in spec.get("slides", []):
        slide = prs.slides.add_slide(layout)

        bg = slide.shapes.add_shape(1, 0, 0, prs.slide_width, prs.slide_height)
        bg.fill.solid(); bg.fill.fore_color.rgb = RGBColor(0xFC, 0xFB, 0xF7); bg.line.fill.background()

        top = slide.shapes.add_shape(1, 0, 0, prs.slide_width, Inches(0.18))
        top.fill.solid(); top.fill.fore_color.rgb = PRIMARY; top.line.fill.background()

        heading_box = slide.shapes.add_textbox(Inches(0.6), Inches(0.45), Inches(12), Inches(0.9))
        _txt(heading_box.text_frame, sd.get("heading", ""), size=28, bold=True, color=PRIMARY)

        bullets = sd.get("bullets", []) or []
        chart = sd.get("chart")

        if chart and chart.get("values"):
            if bullets:
                tb = slide.shapes.add_textbox(Inches(0.6), Inches(1.6), Inches(5.6), Inches(5.4))
                tf = tb.text_frame; tf.word_wrap = True; tf.clear()
                for i, item in enumerate(bullets):
                    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                    p.text = f"•  {item}"
                    for run in p.runs:
                        run.font.size = Pt(16); run.font.color.rgb = TEXT
                    p.space_after = Pt(8)

            png = _chart_png(chart)
            chart_left = Inches(6.6) if bullets else Inches(2.5)
            chart_w = Inches(6.2) if bullets else Inches(8.0)
            slide.shapes.add_picture(io.BytesIO(png), chart_left, Inches(1.6), width=chart_w)
        elif bullets:
            tb = slide.shapes.add_textbox(Inches(0.6), Inches(1.6), Inches(12), Inches(5.4))
            tf = tb.text_frame; tf.word_wrap = True; tf.clear()
            for i, item in enumerate(bullets):
                p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                p.text = f"•  {item}"
                for run in p.runs:
                    run.font.size = Pt(18); run.font.color.rgb = TEXT
                p.space_after = Pt(10)

        if sd.get("notes"):
            slide.notes_slide.notes_text_frame.text = sd["notes"]

    Path(output_path).parent.mkdir(parents=True, exist_ok=True)
    prs.save(output_path)
    return output_path
